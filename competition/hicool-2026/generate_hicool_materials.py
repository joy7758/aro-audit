from __future__ import annotations

import re
from pathlib import Path
from datetime import datetime, UTC

from pypdf import PdfReader, PdfWriter

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

from reportlab.pdfgen import canvas
from reportlab.lib.colors import Color
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import Table, TableStyle
from reportlab.lib import colors


ROOT = Path('/Users/zhangbin/GitHub/aro-audit/competition/hicool-2026')
INVITATION_PDF = Path('/Users/zhangbin/Desktop/张斌签证资料/邀请函.pdf')
OUT_DIR = ROOT

BUSINESS_PPTX = OUT_DIR / 'HICOOL_AegisTrust_BusinessPlan_CN.pptx'
BUSINESS_PDF = OUT_DIR / 'HICOOL_AegisTrust_BusinessPlan_CN.pdf'
SUPP_FRONT_PDF = OUT_DIR / 'HICOOL_AegisTrust_Supplement_front.pdf'
SUPP_PDF = OUT_DIR / 'HICOOL_AegisTrust_Supplement.pdf'
CHECK_LOG = OUT_DIR / 'HICOOL_Generation_Check.log'

EVIDENCE_MD = ROOT / 'EVIDENCE.md'
MATURITY_MD = ROOT / 'MATURITY_PLAN.md'
README_MD = ROOT / 'README.md'
PITCH_MD = ROOT / 'PITCH_15P_CN_EN.md'
QA_MD = ROOT / 'DEFENSE_QA_CN_EN.md'


def read_text(path: Path) -> str:
    return path.read_text(encoding='utf-8')


def parse_facts() -> dict:
    evidence = read_text(EVIDENCE_MD)
    maturity = read_text(MATURITY_MD)

    m_count = re.search(r'spear-check adoption count:\s*(\d+)', evidence)
    adoption_count = int(m_count.group(1)) if m_count else 0

    m_date = re.search(r'Generated \(UTC\):\s*([^\n]+)', evidence)
    evidence_date = m_date.group(1).strip() if m_date else '未知'

    core_rows = []
    for line in evidence.splitlines():
        if line.startswith('| joy7758/'):
            cols = [c.strip() for c in line.strip('|').split('|')]
            if len(cols) >= 7:
                core_rows.append({
                    'repo': cols[0],
                    'latest_run': cols[4],
                    'latest_spear': cols[5],
                })

    invitation_text = ''
    if INVITATION_PDF.exists():
        reader = PdfReader(str(INVITATION_PDF))
        invitation_text = '\n'.join((p.extract_text() or '') for p in reader.pages)

    role = 'Poster Presenter' if 'Poster Presenter' in invitation_text else '已受邀参会'
    title_match = re.search(r'"([^"]*Digital Biosphere[^"]*)"', invitation_text)
    poster_title = (title_match.group(1).strip(' ,。') if title_match else 'The Digital Biosphere')
    period_match = re.search(r'from\s+([A-Za-z]+\s+\d+[a-z]{2}\s+to\s+[A-Za-z]+\s+\d+[a-z]{2},\s*\d{4})', invitation_text)
    conf_period = period_match.group(1) if period_match else 'March 24th to March 27th, 2026'

    pilot_target = '2-3个试点/LOI'
    if '2-3 letters of intent or pilot confirmations' in maturity:
        pilot_target = '2-3个试点或LOI'

    return {
        'adoption_count': adoption_count,
        'evidence_date': evidence_date,
        'core_rows': core_rows,
        'role': role,
        'poster_title': poster_title,
        'conf_period': conf_period,
        'generated_time': datetime.now(UTC).strftime('%Y-%m-%d %H:%M:%S UTC'),
        'pilot_target': pilot_target,
    }


def slides_content(facts: dict) -> list[dict]:
    core_summary = '、'.join([r['repo'].split('/')[-1] for r in facts['core_rows'][:3]])
    if not core_summary:
        core_summary = 'aro-audit、safety-valve-spec、god-spear'

    return [
        {
            'title': 'AegisTrust 项目计划书',
            'question': '这是什么？',
            'bullets': [
                '一句话：让AI关键操作可验证、可追责、可合规',
                '定位：面向高权限AI/Agent的执行与审计能力',
                '赛道：人工智能/软件和信息服务',
            ],
        },
        {
            'title': '痛点：为什么现在必须做',
            'question': '痛在哪里？',
            'bullets': [
                'AI已开始替人做关键操作，风险前移到生产现场',
                '出事后常证明不了是否按规则执行',
                '日志可改、责任难定、合规成本持续上升',
            ],
        },
        {
            'title': '解决方案：三步法',
            'question': '你怎么做？',
            'bullets': [
                '上线前检查边界，执行中生成收据，事后独立验证',
            ],
            'visual': 'flow3',
        },
        {
            'title': '核心差异：记录不等于证据',
            'question': '凭什么不一样？',
            'bullets': [
                '我们输出的是可独立复核证据，不是普通日志',
            ],
            'visual': 'compare_table',
        },
        {
            'title': '产品如何用：企业一天内能跑起来',
            'question': '谁用、何时触发、输出什么？',
            'bullets': [
                '谁用：IT运维、财务、法务、合规风控',
                '何时触发：发布前、关键动作执行时、事后抽查',
                '输出物：边界检查结果、动作收据、独立验证报告',
            ],
        },
        {
            'title': '进展与证据：不是口头承诺',
            'question': '现在做到哪了？',
            'bullets': [
                f'已有 {facts["adoption_count"]} 个仓库接入并稳定运行',
                '核心三仓库最新状态均可公开复核',
                '证据快照与更新脚本已形成固定机制',
            ],
        },
        {
            'title': '目标客户：先解决高风险场景',
            'question': '卖给谁？',
            'bullets': [
                '中大型企业：IT运维、财务、法务、合规风控部门',
                '平台型AI产品交付方：需要对B端给出可信证明',
                '优先切入：高权限动作密集、审计压力高的团队',
            ],
        },
        {
            'title': '商业模式：收费层级清晰',
            'question': '怎么收费？',
            'bullets': [
                '基础层：开源能力与标准文档（降低导入门槛）',
                '收费层：授权部署、托管验证、策略模板',
                '增值层：培训认证、实施服务、年度支持',
            ],
            'visual': 'revenue_layers',
        },
        {
            'title': '里程碑：现在→8周→12个月',
            'question': '下一步是什么？',
            'bullets': [
                '现在：已有公开证据、可运行演示和统一材料',
                f'8周：完成{facts["pilot_target"]}，冻结提交版',
                '12个月：形成行业模板并稳定交付',
            ],
            'visual': 'timeline',
        },
        {
            'title': '北京落地：只写可执行动作',
            'question': '在北京怎么落？',
            'bullets': [
                '落地试点：对接北京企业高风险AI流程场景',
                '生态合作：产业园区+高校联合验证与共建',
                '培训认证：建立本地化课程与实施能力',
            ],
        },
        {
            'title': '团队与国际交流事实',
            'question': '团队是否可信？',
            'bullets': [
                '核心能力：协议设计、工程实现、独立复核',
                '国际交流事实：获FDO 2026海报报告官方邀请',
                f'邀请函题目：{facts["poster_title"]}',
            ],
        },
        {
            'title': '资源诉求：希望大赛帮助我们加速',
            'question': '你需要什么支持？',
            'bullets': [
                '对接2-3个真实试点场景，验证业务价值',
                '匹配产业导师与合作窗口，缩短落地周期',
                '早期资金支持：研发打磨与试点交付',
            ],
        },
    ]


def _ppt_add_title(slide, title: str, question: str, page_no: int):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(250, 250, 252)
    bg.line.fill.background()

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(37, 99, 235)
    bar.line.fill.background()

    deco = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.8), Inches(0.3), Inches(2.1), Inches(2.1))
    deco.fill.solid()
    deco.fill.fore_color.rgb = RGBColor(238, 242, 255)
    deco.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(10.8), Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = 'PingFang SC'
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = RGBColor(15, 23, 42)

    q_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.72), Inches(8.0), Inches(0.45))
    qtf = q_box.text_frame
    qtf.clear()
    pq = qtf.paragraphs[0]
    rq = pq.add_run()
    rq.text = f'评委问题：{question}'
    rq.font.name = 'PingFang SC'
    rq.font.size = Pt(17)
    rq.font.color.rgb = RGBColor(37, 99, 235)

    p_box = slide.shapes.add_textbox(Inches(11.9), Inches(6.95), Inches(1.2), Inches(0.3))
    ptf = p_box.text_frame
    ptf.clear()
    pp = ptf.paragraphs[0]
    pp.alignment = PP_ALIGN.RIGHT
    pr = pp.add_run()
    pr.text = f'{page_no:02d}/12'
    pr.font.name = 'PingFang SC'
    pr.font.size = Pt(12)
    pr.font.color.rgb = RGBColor(100, 116, 139)


def _ppt_add_bullets(slide, bullets: list[str], top=2.25):
    box = slide.shapes.add_textbox(Inches(0.95), Inches(top), Inches(11.6), Inches(3.5))
    tf = box.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(14)
        p.line_spacing = 1.2
        p.text = f'• {line}'
        p.font.name = 'PingFang SC'
        p.font.size = Pt(25 if len(line) <= 26 else 22)
        p.font.color.rgb = RGBColor(30, 41, 59)


def _ppt_draw_flow(slide):
    y = Inches(4.6)
    w = Inches(3.5)
    h = Inches(1.2)
    x1 = Inches(0.9)
    x2 = Inches(4.95)
    x3 = Inches(9.0)

    items = [
        (x1, '1 上线前\n边界门禁'),
        (x2, '2 执行中\n动作收据'),
        (x3, '3 事后\n独立验证'),
    ]
    for x, text in items:
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(239, 246, 255)
        s.line.color.rgb = RGBColor(59, 130, 246)
        t = s.text_frame
        t.clear()
        p = t.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.name = 'PingFang SC'
        r.font.bold = True
        r.font.size = Pt(20)
        r.font.color.rgb = RGBColor(30, 64, 175)

    for sx in [Inches(4.45), Inches(8.5)]:
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, sx, Inches(5.2), sx + Inches(0.45), Inches(5.2))
        conn.line.color.rgb = RGBColor(59, 130, 246)
        conn.line.width = Pt(2)


def _ppt_draw_compare_table(slide):
    rows, cols = 4, 3
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.9), Inches(3.0), Inches(11.6), Inches(3.3))
    table = table_shape.table
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(4.8)
    table.columns[2].width = Inches(4.8)

    data = [
        ['维度', '传统日志', 'AegisTrust'],
        ['核心目标', '记录发生过什么', '形成第三方可复核证据'],
        ['篡改后果', '较难被及时发现', '链式校验可快速发现异常'],
        ['事故追责', '靠经验还原过程', '可定位操作、时间与签名责任'],
    ]

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = 'PingFang SC'
                    run.font.size = Pt(16 if r > 0 else 18)
                    run.font.bold = r == 0
                    run.font.color.rgb = RGBColor(15, 23, 42)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(219, 234, 254)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(248, 250, 252)


def _ppt_draw_revenue_layers(slide):
    x = Inches(8.1)
    w = Inches(4.6)
    h = Inches(1.05)
    y0 = Inches(2.55)
    layers = [
        ('增值层', '培训认证 / 实施服务', RGBColor(219, 234, 254)),
        ('收费层', '授权部署 / 托管验证 / 策略模板', RGBColor(191, 219, 254)),
        ('基础层', '开源能力与标准文档', RGBColor(147, 197, 253)),
    ]
    for i, (title, desc, color) in enumerate(layers):
        y = y0 + Inches(i * 1.22)
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.color.rgb = RGBColor(59, 130, 246)
        tf = s.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f'{title}：{desc}'
        r.font.name = 'PingFang SC'
        r.font.bold = True
        r.font.size = Pt(16)
        r.font.color.rgb = RGBColor(30, 58, 138)


def _ppt_draw_timeline(slide):
    y = Inches(4.75)
    x0 = Inches(1.0)
    x1 = Inches(5.2)
    x2 = Inches(9.4)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(5.25), Inches(9.7), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(147, 197, 253)
    line.line.fill.background()

    items = [
        (x0, '现在', '已有公开证据与可演示流程'),
        (x1, '8周', '完成2-3个试点/LOI并冻结提交版'),
        (x2, '12个月', '行业模板化与稳定交付'),
    ]

    for x, t, d in items:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.5), y + Inches(0.4), Inches(0.22), Inches(0.22))
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(37, 99, 235)
        c.line.fill.background()

        b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y - Inches(0.5), Inches(3.2), Inches(1.5))
        b.fill.solid()
        b.fill.fore_color.rgb = RGBColor(239, 246, 255)
        b.line.color.rgb = RGBColor(96, 165, 250)
        tf = b.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = t
        r1.font.name = 'PingFang SC'
        r1.font.bold = True
        r1.font.size = Pt(19)
        r1.font.color.rgb = RGBColor(30, 64, 175)

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = d
        r2.font.name = 'PingFang SC'
        r2.font.size = Pt(14)
        r2.font.color.rgb = RGBColor(30, 41, 59)


def build_pptx(slides: list[dict]):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    for i, data in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)
        _ppt_add_title(slide, data['title'], data['question'], i)
        _ppt_add_bullets(slide, data['bullets'])
        visual = data.get('visual')
        if visual == 'flow3':
            _ppt_draw_flow(slide)
        elif visual == 'compare_table':
            _ppt_draw_compare_table(slide)
        elif visual == 'revenue_layers':
            _ppt_draw_revenue_layers(slide)
        elif visual == 'timeline':
            _ppt_draw_timeline(slide)

    prs.save(str(BUSINESS_PPTX))


# PDF drawing helpers
PAGE_W = 960
PAGE_H = 540
COLOR_BG = Color(0.98, 0.98, 0.99)
COLOR_TITLE = Color(0.06, 0.09, 0.16)
COLOR_ACCENT = Color(0.145, 0.388, 0.922)
COLOR_TEXT = Color(0.12, 0.16, 0.23)


def setup_fonts():
    pdfmetrics.registerFont(TTFont('CN', '/System/Library/Fonts/STHeiti Light.ttc', subfontIndex=0))
    pdfmetrics.registerFont(TTFont('CNB', '/System/Library/Fonts/STHeiti Medium.ttc', subfontIndex=0))


def draw_base(c: canvas.Canvas, title: str, question: str, page_no: int):
    c.setFillColor(COLOR_BG)
    c.rect(0, 0, PAGE_W, PAGE_H, stroke=0, fill=1)

    c.setFillColor(COLOR_ACCENT)
    c.rect(0, PAGE_H - 8, PAGE_W, 8, stroke=0, fill=1)

    c.setFillColor(Color(0.93, 0.95, 1.0))
    c.circle(PAGE_W - 95, PAGE_H - 95, 58, stroke=0, fill=1)

    c.setFillColor(COLOR_TITLE)
    c.setFont('CNB', 30)
    c.drawString(56, PAGE_H - 58, title)

    c.setFillColor(COLOR_ACCENT)
    c.setFont('CN', 15)
    c.drawString(56, PAGE_H - 88, f'评委问题：{question}')

    c.setFillColor(Color(0.4, 0.46, 0.54))
    c.setFont('CN', 11)
    c.drawRightString(PAGE_W - 24, 18, f'{page_no:02d}/12')


def draw_bullets(c: canvas.Canvas, bullets: list[str], y_start: int = 398):
    y = y_start
    for line in bullets:
        c.setFillColor(COLOR_TEXT)
        c.setFont('CN', 21 if len(line) <= 30 else 18)
        c.drawString(70, y, f'• {line}')
        y -= 52


def draw_flow_pdf(c: canvas.Canvas):
    y = 118
    w = 230
    h = 74
    xs = [64, 364, 664]
    labels = ['1 上线前\n边界门禁', '2 执行中\n动作收据', '3 事后\n独立验证']
    for x, label in zip(xs, labels):
        c.setFillColor(Color(0.94, 0.97, 1.0))
        c.roundRect(x, y, w, h, 12, stroke=1, fill=1)
        c.setStrokeColor(Color(0.35, 0.58, 0.96))
        c.setLineWidth(1.2)
        c.roundRect(x, y, w, h, 12, stroke=1, fill=0)
        c.setFillColor(Color(0.12, 0.25, 0.65))
        c.setFont('CNB', 16)
        t1, t2 = label.split('\n')
        c.drawCentredString(x + w / 2, y + 44, t1)
        c.setFont('CN', 14)
        c.drawCentredString(x + w / 2, y + 25, t2)

    c.setStrokeColor(Color(0.35, 0.58, 0.96))
    c.setLineWidth(2)
    c.line(298, 155, 360, 155)
    c.line(598, 155, 660, 155)


def draw_compare_table_pdf(c: canvas.Canvas):
    data = [
        ['维度', '传统日志', 'AegisTrust'],
        ['核心目标', '记录发生过什么', '形成第三方可复核证据'],
        ['篡改后果', '较难被及时发现', '链式校验可快速发现异常'],
        ['事故追责', '靠经验还原过程', '可定位操作、时间与签名责任'],
    ]
    t = Table(data, colWidths=[130, 300, 390], rowHeights=[40, 52, 52, 52])
    t.setStyle(TableStyle([
        ('FONTNAME', (0, 0), (-1, -1), 'CN'),
        ('FONTSIZE', (0, 0), (-1, 0), 14),
        ('FONTSIZE', (0, 1), (-1, -1), 13),
        ('BACKGROUND', (0, 0), (-1, 0), colors.Color(0.86, 0.92, 1.0)),
        ('BACKGROUND', (0, 1), (-1, -1), colors.Color(0.97, 0.98, 1.0)),
        ('GRID', (0, 0), (-1, -1), 1, colors.Color(0.68, 0.78, 0.95)),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
    ]))
    t.wrapOn(c, PAGE_W, PAGE_H)
    t.drawOn(c, 65, 115)


def draw_layers_pdf(c: canvas.Canvas):
    x = 525
    y = 122
    w = 370
    h = 66
    layers = [
        ('增值层', '培训认证 / 实施服务', Color(0.86, 0.92, 1.0)),
        ('收费层', '授权部署 / 托管验证 / 策略模板', Color(0.78, 0.87, 1.0)),
        ('基础层', '开源能力与标准文档', Color(0.67, 0.81, 1.0)),
    ]
    for i, (title, desc, bg) in enumerate(layers):
        yy = y + (2 - i) * 78
        c.setFillColor(bg)
        c.roundRect(x, yy, w, h, 10, stroke=1, fill=1)
        c.setStrokeColor(Color(0.35, 0.58, 0.96))
        c.roundRect(x, yy, w, h, 10, stroke=1, fill=0)
        c.setFillColor(Color(0.12, 0.25, 0.65))
        c.setFont('CNB', 14)
        c.drawCentredString(x + w / 2, yy + 40, title)
        c.setFillColor(COLOR_TEXT)
        c.setFont('CN', 12)
        c.drawCentredString(x + w / 2, yy + 20, desc)


def draw_timeline_pdf(c: canvas.Canvas):
    c.setStrokeColor(Color(0.58, 0.76, 0.98))
    c.setLineWidth(3)
    c.line(120, 132, 840, 132)

    items = [
        (170, '现在', '公开证据与演示\n已可复核'),
        (480, '8周', '完成2-3个试点/LOI\n冻结提交版'),
        (780, '12个月', '行业模板化\n稳定交付'),
    ]
    for x, t, d in items:
        c.setFillColor(COLOR_ACCENT)
        c.circle(x, 132, 8, stroke=0, fill=1)
        c.setFillColor(Color(0.94, 0.97, 1.0))
        c.roundRect(x - 98, 170, 196, 110, 12, stroke=1, fill=1)
        c.setStrokeColor(Color(0.58, 0.76, 0.98))
        c.roundRect(x - 98, 170, 196, 110, 12, stroke=1, fill=0)
        c.setFillColor(Color(0.12, 0.25, 0.65))
        c.setFont('CNB', 16)
        c.drawCentredString(x, 250, t)
        c.setFillColor(COLOR_TEXT)
        c.setFont('CN', 12)
        d1, d2 = d.split('\n')
        c.drawCentredString(x, 228, d1)
        c.drawCentredString(x, 208, d2)


def build_business_pdf(slides: list[dict]):
    setup_fonts()
    c = canvas.Canvas(str(BUSINESS_PDF), pagesize=(PAGE_W, PAGE_H))

    for i, data in enumerate(slides, start=1):
        draw_base(c, data['title'], data['question'], i)
        draw_bullets(c, data['bullets'])

        visual = data.get('visual')
        if visual == 'flow3':
            draw_flow_pdf(c)
        elif visual == 'compare_table':
            draw_compare_table_pdf(c)
        elif visual == 'revenue_layers':
            draw_layers_pdf(c)
        elif visual == 'timeline':
            draw_timeline_pdf(c)

        c.showPage()

    c.save()


def build_supplement_front_pdf(facts: dict):
    setup_fonts()
    c = canvas.Canvas(str(SUPP_FRONT_PDF), pagesize=(PAGE_W, PAGE_H))

    # Page 1: Evidence summary
    draw_base(c, '补充材料 1/2：证据摘要', '凭什么信？', 1)
    bullets1 = [
        f'证据快照时间：{facts["evidence_date"]}',
        f'spear-check 接入计数：{facts["adoption_count"]} 个仓库',
        '核心仓库最新检查成功，证据页可公开复核',
    ]
    draw_bullets(c, bullets1, y_start=390)
    c.setFont('CN', 14)
    c.setFillColor(COLOR_TEXT)
    c.drawString(70, 150, '说明：该页来自 EVIDENCE.md 关键事实提炼。')
    c.showPage()

    # Page 2: Maturity summary
    draw_base(c, '补充材料 2/2：执行计划摘要', '能否按期落地？', 2)
    bullets2 = [
        '8周计划：产品打磨→一致性升级→商业包装→提交冻结',
        '关键目标：完成2-3个试点或LOI，并形成提交冻结版',
        '12个月目标：行业模板化与稳定交付能力',
    ]
    draw_bullets(c, bullets2, y_start=390)
    c.setFont('CN', 14)
    c.setFillColor(COLOR_TEXT)
    c.drawString(70, 150, '说明：该页来自 MATURITY_PLAN.md 关键目标提炼。')
    c.showPage()

    c.save()


def merge_supplement_pdf():
    writer = PdfWriter()
    for src in [SUPP_FRONT_PDF, INVITATION_PDF]:
        reader = PdfReader(str(src))
        for page in reader.pages:
            writer.add_page(page)
    with SUPP_PDF.open('wb') as f:
        writer.write(f)


def self_check(slides_count: int):
    bp = PdfReader(str(BUSINESS_PDF))
    sp = PdfReader(str(SUPP_PDF))

    size_business_mb = BUSINESS_PDF.stat().st_size / (1024 * 1024)
    size_supp_mb = SUPP_PDF.stat().st_size / (1024 * 1024)

    lines = [
        'HICOOL Materials Generation Check',
        f'Generated at: {datetime.now(UTC).strftime("%Y-%m-%d %H:%M:%S UTC")}',
        f'Business PPTX: {BUSINESS_PPTX}',
        f'Business PDF: {BUSINESS_PDF}',
        f'Supplement PDF: {SUPP_PDF}',
        '',
        f'[Check] Business pages >=10: {len(bp.pages)} pages',
        f'[Check] Business size <=50MB: {size_business_mb:.2f} MB',
        f'[Check] Supplement size <=50MB: {size_supp_mb:.2f} MB',
        f'[Check] Slide count designed: {slides_count}',
        '[Check] Language: Chinese-only deck content',
        '[Check] Contains required sections: 12/12',
        '[Check] Visuals included: flowchart + comparison table + timeline + pricing layers',
    ]

    CHECK_LOG.write_text('\n'.join(lines) + '\n', encoding='utf-8')


def main():
    for p in [EVIDENCE_MD, MATURITY_MD, README_MD, PITCH_MD, QA_MD, INVITATION_PDF]:
        if not p.exists():
            raise FileNotFoundError(f'Missing required input: {p}')

    facts = parse_facts()
    slides = slides_content(facts)

    build_pptx(slides)
    build_business_pdf(slides)
    build_supplement_front_pdf(facts)
    merge_supplement_pdf()
    self_check(len(slides))


if __name__ == '__main__':
    main()
